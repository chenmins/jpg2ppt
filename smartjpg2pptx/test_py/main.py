import cv2
import numpy as np
import os
import torch
from pptx import Presentation
from pptx.util import Inches
from glob import glob
import detectron2
from detectron2.utils.logger import setup_logger
from detectron2 import model_zoo
from detectron2.engine import DefaultPredictor
from detectron2.config import get_cfg
from detectron2.utils.visualizer import Visualizer
import matplotlib.pyplot as plt

# 设置Detectron2的日志记录
setup_logger()
print(torch.cuda.is_available())

def setup_cfg():
    # 加载配置和预训练模型
    cfg = get_cfg()
    cfg.merge_from_file(model_zoo.get_config_file("COCO-InstanceSegmentation/mask_rcnn_R_50_FPN_3x.yaml"))
    cfg.MODEL.ROI_HEADS.SCORE_THRESH_TEST = 0.3  # 设置阈值
    cfg.MODEL.WEIGHTS = model_zoo.get_checkpoint_url("COCO-InstanceSegmentation/mask_rcnn_R_50_FPN_3x.yaml")
    cfg.MODEL.DEVICE = "cuda" if torch.cuda.is_available() else "cpu"  # 使用GPU如果可用
    return cfg


def find_screen_contour(image_path, cfg):
    predictor = DefaultPredictor(cfg)
    img = cv2.imread(image_path)
    outputs = predictor(img)

    # 使用正确的导入创建Visualizer实例
    v = Visualizer(img[:, :, ::-1], metadata=detectron2.data.MetadataCatalog.get(cfg.DATASETS.TRAIN[0]), scale=0.5)
    v = v.draw_instance_predictions(outputs["instances"].to("cpu"))
    plt.imshow(v.get_image()[:, :, ::-1])
    plt.show()

    # screen_id = 62  # TV
    # screen_id = 73  # BOOK
    screen_id = [62,73]
    found = False

    for i, (cls, score) in enumerate(zip(outputs["instances"].pred_classes, outputs["instances"].scores)):
        if cls in screen_id and score > cfg.MODEL.ROI_HEADS.SCORE_THRESH_TEST:
            mask = outputs["instances"].pred_masks[i].cpu().numpy()
            y, x = np.where(mask)
            if y.size > 0 and x.size > 0:
                cropped_img = img[min(y):max(y), min(x):max(x)]
                found = True
                print(f"Screen found in {image_path}, cropping...")

                # 修正梯形效果的透视变换
                cropped_img = perspective_correction(cropped_img)

                return cropped_img

    if not found:
        print(f"No screen found in {image_path}.")
    return None


def perspective_correction(image):
    # 通过透视变换修正梯形效果
    # 假设你已经有了一个四点坐标来校正梯形
    height, width = image.shape[:2]
    src_points = np.float32([[0, 0], [width, 0], [width, height], [0, height]])  # 原图四个角点
    dst_points = np.float32(
        [[100, 50], [width - 100, 50], [width - 100, height - 50], [100, height - 50]])  # 目标四个角点（校正后）

    # 计算透视变换矩阵
    matrix = cv2.getPerspectiveTransform(src_points, dst_points)

    # 应用透视变换
    corrected_img = cv2.warpPerspective(image, matrix, (width, height))
    return corrected_img


def create_presentation(images, cfg, output_file='output_presentation.pptx'):
    prs = Presentation()
    slide_width = Inches(10)  # 设置幻灯片的宽度
    slide_height = Inches(7.5)  # 设置幻灯片的高度

    for image_path in images:
        cropped_image = find_screen_contour(image_path, cfg)
        if cropped_image is not None:
            # 保存处理后的图片
            cropped_image_path = f"cropped_{os.path.basename(image_path)}"
            cv2.imwrite(cropped_image_path, cropped_image)
            # 创建新的幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # 在幻灯片中添加图片，调整为全屏显示
            slide.shapes.add_picture(cropped_image_path, 0, 0, width=slide_width, height=slide_height)
    # 保存PPTX文件
    prs.save(output_file)


def main():
    folder_path = 'd:/test1'
    images = sorted(glob(os.path.join(folder_path, '*.jpg')))
    cfg = setup_cfg()
    create_presentation(images, cfg)


if __name__ == "__main__":
    main()
