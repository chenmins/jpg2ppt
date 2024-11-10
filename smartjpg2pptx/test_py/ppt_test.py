from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Hello, PowerPoint!"

# 添加一张图片，确保这张图片是可以访问的
img_path = 'D:/培训课件/01国务院国资委AI+专项行动相关政策解读/20241104_091635.jpg'
left = Inches(2)
top = Inches(2)
slide.shapes.add_picture(img_path, left, top, width=Inches(4), height=Inches(3))

prs.save('test_presentation.pptx')
