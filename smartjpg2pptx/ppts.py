import cv2
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches
from glob import glob

def add_image_to_slide(prs, image_path, slide_width, slide_height):
    img = cv2.imdecode(np.fromfile(image_path, dtype=np.uint8), cv2.IMREAD_UNCHANGED)
    if img is None:
        print(f"Warning: Unable to read image at {image_path}. Skipping...")
        return

    h, w, _ = img.shape
    ratio = min(slide_width / w, slide_height / h)
    img_width = int(w * ratio)
    img_height = int(h * ratio)
    left = int((slide_width - img_width) / 2)
    top = int((slide_height - img_height) / 2)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(image_path, left, top, width=img_width, height=img_height)

def create_presentation(images, output_file):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(10 * 9 / 16)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for image_path in images:
        add_image_to_slide(prs, image_path, slide_width, slide_height)

    prs.save(output_file)

def main():
    folder_path = 'D:/培训课件'
    for root, dirs, files in os.walk(folder_path):
        for dir_name in dirs:
            sub_folder_path = os.path.join(root, dir_name)
            images = sorted(glob(os.path.join(sub_folder_path, '*.jpg')))
            if images:
                output_file = os.path.join(folder_path, f"{dir_name}.pptx")
                create_presentation(images, output_file)

if __name__ == "__main__":
    main()
